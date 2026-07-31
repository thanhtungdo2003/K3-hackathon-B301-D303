"""Đọc file PPTX/PDF thật thành các block để frontend vẽ lên HTML canvas.

PPTX: không render ảnh slide — lấy phần chữ có cấu trúc (tiêu đề, gạch đầu dòng,
bảng, ghi chú của người trình bày) rồi vẽ lại bằng canvas. Nhờ vậy nội dung
slide vẫn là văn bản: đọc được, tìm được, và Advisor biết được tiêu đề slide.

PDF: render thêm ảnh từng trang đúng như bản gốc để hiển thị, nhưng **vẫn** trích
văn bản thành block. Ảnh dùng để xem, block dùng làm ngữ cảnh cho Advisor.
"""
from __future__ import annotations

import uuid
from pathlib import Path
from typing import Any

import fitz  # pymupdf — đọc chữ và render ảnh cho PDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

MAX_BULLETS = 8
MAX_CHARS = 220

# Bề ngang ảnh trang PDF. 1600px đủ nét cho máy chiếu mà file vẫn nhẹ.
PAGE_IMAGE_WIDTH = 1600

# Nơi ảnh trang PDF được phục vụ (khớp với mount trong main.py).
PAGE_URL_PREFIX = "/slide-pages"


def page_image_url(name: str | None) -> str | None:
    return f"{PAGE_URL_PREFIX}/{name}" if name else None


def _clean(text: str) -> str:
    return " ".join(text.replace("\x0b", " ").split()).strip()


def _looks_like_code(lines: list[str]) -> bool:
    """Đoạn nhiều ký hiệu lập trình thì render bằng font mono."""
    if len(lines) < 2:
        return False
    hits = sum(1 for ln in lines if any(tok in ln for tok in ("()", "{", "};", "def ", "class ", "=>", "->", "import ")))
    return hits >= max(2, len(lines) // 2)


def _shape_text_lines(shape: Any) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = _clean("".join(run.text for run in para.runs) or para.text)
        if text:
            lines.append(text)
    return lines


def parse_pptx(path: Path) -> list[dict]:
    """Trả về danh sách slide: {index, title, blocks, notes}."""
    prs = Presentation(str(path))
    result: list[dict] = []

    for i, slide in enumerate(prs.slides):
        title = ""
        blocks: list[dict] = []
        body_chunks: list[list[str]] = []
        tables: list[list[list[str]]] = []
        picture_count = 0

        # Tiêu đề lấy từ placeholder title nếu có
        try:
            if slide.shapes.title is not None:
                title = _clean(slide.shapes.title.text)
        except (AttributeError, KeyError):
            title = ""

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                continue

            if getattr(shape, "has_table", False):
                rows = [
                    [_clean(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    tables.append(rows)
                continue

            if shape is getattr(slide.shapes, "title", None):
                continue

            lines = _shape_text_lines(shape)
            if lines:
                body_chunks.append(lines)

        # Không có placeholder title -> lấy dòng đầu tiên làm tiêu đề
        if not title and body_chunks:
            first = body_chunks[0]
            title = first[0][:200]
            rest = first[1:]
            body_chunks = ([rest] if rest else []) + body_chunks[1:]

        if title:
            blocks.append({"type": "title", "text": title[:200]})

        for chunk in body_chunks:
            trimmed = [ln[:MAX_CHARS] for ln in chunk][:MAX_BULLETS]
            if not trimmed:
                continue
            if _looks_like_code(trimmed):
                blocks.append({"type": "code", "lines": trimmed})
            elif len(trimmed) == 1 and len(trimmed[0]) <= 140:
                blocks.append({"type": "lead", "text": trimmed[0]})
            else:
                blocks.append({"type": "bullets", "items": trimmed})

        for rows in tables:
            blocks.append({"type": "table", "rows": [r[:4] for r in rows[:6]]})

        if picture_count:
            blocks.append(
                {"type": "note", "text": f"Slide gốc có {picture_count} hình ảnh — xem bản PPTX để đối chiếu."}
            )

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = _clean(slide.notes_slide.notes_text_frame.text)

        result.append(
            {
                "index": i,
                "title": title or f"Slide {i + 1}",
                "blocks": blocks or [{"type": "title", "text": f"Slide {i + 1}"}],
                "notes": notes[:2000],
            }
        )

    return result


def _render_page(page: Any, out_dir: Path, name: str) -> str:
    """Render một trang thành PNG trong `out_dir`. Trả về tên file."""
    # Phóng theo bề ngang trang để mọi trang ra cùng độ nét, không phụ thuộc khổ
    # giấy gốc (A4 dọc, 16:9 ngang…).
    zoom = PAGE_IMAGE_WIDTH / max(page.rect.width, 1)
    pixmap = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    pixmap.save(str(out_dir / name))
    return name


def parse_pdf(path: Path, page_image_dir: Path | None = None) -> list[dict]:
    """Trả về danh sách slide: {index, title, blocks, notes, page_image}.

    Dùng PyMuPDF cho cả chữ lẫn ảnh. Không dùng pypdf vì pypdf đòi thêm gói
    `cryptography` mới đọc nổi PDF đã mã hoá (rất nhiều PDF xuất ra có đặt hạn
    chế quyền), còn PyMuPDF mở thẳng được.

    `page_image_dir` có giá trị thì render thêm ảnh từng trang vào đó; render
    hỏng thì bỏ ảnh và vẫn trả về phần văn bản, không làm hỏng cả lần tải lên.
    """
    result: list[dict] = []
    stem = uuid.uuid4().hex[:10]
    if page_image_dir is not None:
        page_image_dir.mkdir(parents=True, exist_ok=True)

    with fitz.open(str(path)) as doc:
        if doc.needs_pass:
            raise ValueError("PDF này có mật khẩu mở file.")
        pages = list(doc)
        for i, page in enumerate(pages):
            image = ""
            if page_image_dir is not None:
                try:
                    image = _render_page(page, page_image_dir, f"{stem}-{i + 1:03d}.png")
                except Exception:  # noqa: BLE001 — ảnh là phần tăng thêm, không được chặn import
                    image = ""

            text = page.get_text() or ""
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            title = lines[0][:200] if lines else f"Trang {i + 1}"
            body = lines[1:] if lines else []
            blocks: list[dict] = []

            if title:
                blocks.append({"type": "title", "text": title})

            if body:
                if len(body) == 1 and len(body[0]) <= 140:
                    blocks.append({"type": "lead", "text": body[0]})
                else:
                    items = [line[:MAX_CHARS] for line in body][:MAX_BULLETS]
                    blocks.append({"type": "bullets", "items": items})
            else:
                blocks.append({"type": "note", "text": "Trang PDF này không có văn bản có thể trích xuất."})

            result.append(
                {
                    "index": i,
                    "title": title,
                    "blocks": blocks,
                    "notes": "",
                    "page_image": image,
                }
            )

    return result


def slide_plain_text(slide_row: Any) -> str:
    """Gộp nội dung một slide thành văn bản thuần — dùng làm ngữ cảnh cho LLM."""
    parts: list[str] = [slide_row.title]
    for block in slide_row.blocks or []:
        kind = block.get("type")
        if kind in ("title", "lead", "note", "kicker"):
            parts.append(block.get("text", ""))
        elif kind == "bullets":
            parts.extend(block.get("items", []))
        elif kind == "code":
            parts.extend(block.get("lines", []))
        elif kind == "table":
            for row in block.get("rows", []):
                parts.append(" | ".join(row))
    if getattr(slide_row, "notes", ""):
        parts.append(f"Ghi chú của giảng viên: {slide_row.notes}")
    return "\n".join(p for p in parts if p)[:3000]
